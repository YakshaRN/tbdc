"""
Document Extraction Service.

Extracts text content from various document types including:
- PDF files
- Word documents (.doc, .docx)
- PowerPoint presentations (.ppt, .pptx)
- Text files (.txt, .rtf)
- Excel spreadsheets (.xls, .xlsx)
"""
import io
from typing import Optional, Dict, Any, List
from loguru import logger


class DocumentExtractor:
    """
    Service to extract text content from various document types.
    
    Uses different libraries based on file type:
    - PyPDF2 for PDFs
    - python-docx for Word documents
    - python-pptx for PowerPoint files
    - openpyxl for Excel files
    """
    
    # Maximum text length to extract (to avoid huge prompts)
    MAX_TEXT_LENGTH = 15000
    
    def extract_text(self, content: bytes, file_name: str) -> Optional[str]:
        """
        Extract text from document content.
        
        Args:
            content: File content as bytes
            file_name: Name of the file (used to determine type)
            
        Returns:
            Extracted text or None if extraction failed
        """
        if not content:
            return None
        
        ext = self._get_extension(file_name)
        
        try:
            if ext == '.pdf':
                return self._extract_pdf(content)
            elif ext in ['.doc', '.docx']:
                return self._extract_word(content)
            elif ext in ['.ppt', '.pptx']:
                return self._extract_powerpoint(content)
            elif ext in ['.xls', '.xlsx']:
                return self._extract_excel(content)
            elif ext in ['.txt', '.rtf']:
                return self._extract_text(content)
            else:
                logger.warning(f"Unsupported file type: {ext}")
                return None
        except Exception as e:
            logger.error(f"Error extracting text from {file_name}: {e}")
            return None
    
    def extract_from_attachments(
        self, 
        attachments: List[Dict[str, Any]]
    ) -> Dict[str, str]:
        """
        Extract text from multiple attachments.
        
        Args:
            attachments: List of attachment dicts with 'content' and 'file_name'
            
        Returns:
            Dict mapping file_name to extracted text
        """
        results = {}
        
        for attachment in attachments:
            file_name = attachment.get("file_name", "unknown")
            content = attachment.get("content")
            
            if content:
                text = self.extract_text(content, file_name)
                if text:
                    # Truncate if too long
                    if len(text) > self.MAX_TEXT_LENGTH:
                        text = text[:self.MAX_TEXT_LENGTH] + "\n\n[... content truncated ...]"
                    results[file_name] = text
                    logger.info(f"Extracted {len(text)} chars from {file_name}")
        
        return results
    
    def combine_extracted_text(
        self, 
        extractions: Dict[str, str],
        max_total_length: int = 30000
    ) -> str:
        """
        Combine extracted text from multiple documents into a single string.
        
        Args:
            extractions: Dict mapping file_name to text
            max_total_length: Maximum total text length
            
        Returns:
            Combined text with document headers
        """
        if not extractions:
            return ""
        
        parts = []
        total_length = 0
        
        for file_name, text in extractions.items():
            header = f"\n\n--- Content from: {file_name} ---\n\n"
            
            # Check if we'd exceed the limit
            new_length = total_length + len(header) + len(text)
            if new_length > max_total_length:
                # Truncate this document
                available = max_total_length - total_length - len(header) - 50
                if available > 200:
                    text = text[:available] + "\n[... truncated ...]"
                    parts.append(header + text)
                break
            
            parts.append(header + text)
            total_length = new_length
        
        return "".join(parts)
    
    def _get_extension(self, file_name: str) -> str:
        """Get lowercase file extension."""
        if "." not in file_name:
            return ""
        return "." + file_name.split(".")[-1].lower()
    
    def _extract_pdf(self, content: bytes) -> Optional[str]:
        """Extract text from PDF."""
        try:
            import pypdf
            
            reader = pypdf.PdfReader(io.BytesIO(content))
            text_parts = []
            
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            
            return "\n\n".join(text_parts) if text_parts else None
            
        except ImportError:
            logger.error("pypdf not installed. Install with: pip install pypdf")
            return None
        except Exception as e:
            logger.error(f"Error extracting PDF: {e}")
            return None
    
    def _extract_word(self, content: bytes) -> Optional[str]:
        """Extract text from Word documents."""
        try:
            from docx import Document
            
            doc = Document(io.BytesIO(content))
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        text_parts.append(row_text)
            
            return "\n".join(text_parts) if text_parts else None
            
        except ImportError:
            logger.error("python-docx not installed. Install with: pip install python-docx")
            return None
        except Exception as e:
            logger.error(f"Error extracting Word doc: {e}")
            return None
    
    def _extract_powerpoint(self, content: bytes) -> Optional[str]:
        """Extract text from PowerPoint presentations."""
        try:
            from pptx import Presentation
            
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"[Slide {slide_num}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # Has content beyond slide number
                    text_parts.append("\n".join(slide_text))
            
            return "\n\n".join(text_parts) if text_parts else None
            
        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return None
        except Exception as e:
            logger.error(f"Error extracting PowerPoint: {e}")
            return None
    
    def _extract_excel(self, content: bytes) -> Optional[str]:
        """Extract text from Excel spreadsheets."""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"[Sheet: {sheet_name}]"]
                
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        sheet_text.append(" | ".join(row_values))
                
                if len(sheet_text) > 1:
                    text_parts.append("\n".join(sheet_text))
            
            return "\n\n".join(text_parts) if text_parts else None
            
        except ImportError:
            logger.error("openpyxl not installed. Install with: pip install openpyxl")
            return None
        except Exception as e:
            logger.error(f"Error extracting Excel: {e}")
            return None
    
    def _extract_text(self, content: bytes) -> Optional[str]:
        """Extract text from plain text files."""
        try:
            # Try UTF-8 first, fall back to latin-1
            try:
                return content.decode('utf-8')
            except UnicodeDecodeError:
                return content.decode('latin-1')
        except Exception as e:
            logger.error(f"Error extracting text file: {e}")
            return None


# Global instance
document_extractor = DocumentExtractor()
